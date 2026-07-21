"""PPTX extraction — slide text and speaker notes."""

from __future__ import annotations

from io import BytesIO

from apps.attachments.services.handlers import ExtractionResult, metadata_only


def extract(data: bytes, *, filename: str = "", limit: int = 400_000) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return metadata_only("pptx", "pptx reader unavailable")

    try:
        presentation = Presentation(BytesIO(data))
        parts: list[str] = []
        for number, slide in enumerate(presentation.slides, start=1):
            parts.append(f"[slide {number}]")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
            notes = getattr(slide, "notes_slide", None)
            if notes is not None and getattr(notes, "notes_text_frame", None) is not None:
                note_text = notes.notes_text_frame.text.strip()
                if note_text:
                    parts.append(f"[notes] {note_text}")
        text = "\n".join(parts).strip()
        if not text:
            return metadata_only("pptx", "no text content")
        return ExtractionResult(
            text=text[:limit], page_count=len(presentation.slides), handler="pptx",
            truncated=len(text) > limit,
        )
    except Exception as exc:  # noqa: BLE001
        return metadata_only("pptx", f"unreadable: {exc}"[:200])
